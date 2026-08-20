# -*- coding: utf-8 -*-
"""
PPTX Description 패널(desc-table) 검증 스크립트.

용도: 화면설계서_부품공정자동화현황_샘플.html의 "PPTX 다운로드" 버튼으로 생성한
.pptx 파일을 python-pptx로 열어, 텍스트 잘림/배경 불일치/겹침을 정적으로 검출한다.
사람이 PowerPoint를 열어 48개 슬라이드를 눈으로 확인하지 않고도 회귀를 빠르게
잡기 위한 도구다. (배경: 2026-08-20 desc-table 잘림/겹침 근본 수정 작업 중 만들어짐.
근본 원인과 수정 구조는 memory의 feedback-pptx-desc-table-flow-fix 참고.)

사용법:
    python pptx_desc_검증.py <다운로드한.pptx>
    (인자 생략 시 기본값 downloaded.pptx)

필요 패키지: python-pptx, Pillow (pip install python-pptx Pillow)
필요 폰트: C:/Windows/Fonts/malgun.ttf, malgunbd.ttf (Windows 기본 내장)

6가지 검증을 순서대로 실행하고, 모두 "Total findings: 0"이면 통과.
하나라도 findings가 있으면 실제 PPTX를 열어 해당 슬라이드를 육안으로 재확인할 것
(이 스크립트는 근사 계산이라 완전한 대체가 아님 — 특히 폭 계산에 쓰는 문자 단위
그리디 wrap 시뮬레이션은 PowerPoint의 실제 줄바꿈 알고리즘과 100% 동일하지 않다).
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from PIL import ImageFont

SLIDE_W_IN = 13.33
SLIDE_H_IN = 7.5
EMU_PER_IN = 914400
DPI = 96.0
DESC_PANEL_LEFT_RATIO = 0.68  # Description 패널은 슬라이드 우측 ~24%에 위치

FONT_REG = "C:/Windows/Fonts/malgun.ttf"
FONT_BOLD = "C:/Windows/Fonts/malgunbd.ttf"

_font_cache = {}


def get_font(pt, bold):
    key = (round(pt, 1), bold)
    if key not in _font_cache:
        path = FONT_BOLD if bold else FONT_REG
        px_size = max(6, int(round(pt * DPI / 72.0)))
        _font_cache[key] = ImageFont.truetype(path, px_size)
    return _font_cache[key]


def text_width_in(text, pt, bold):
    if not text:
        return 0.0
    f = get_font(pt, bold)
    bbox = f.getbbox(text)
    return (bbox[2] - bbox[0]) / DPI


def emu_to_in(v):
    return v / EMU_PER_IN if v is not None else 0.0


def wrap_lines(text, pt, bold, avail_w_in):
    """문자 단위 그리디 wrap 시뮬레이션(맑은 고딕 실측 폭 기준)."""
    if avail_w_in <= 0.01:
        return max(1, len(text.split('\n')))
    total_lines = 0
    for manual_line in text.split('\n'):
        if manual_line == '':
            total_lines += 1
            continue
        cur = ''
        lines_this = 1
        for ch in manual_line:
            test = cur + ch
            if text_width_in(test, pt, bold) <= avail_w_in or cur == '':
                cur = test
            else:
                lines_this += 1
                cur = ch
        total_lines += lines_this
    return total_lines


def get_fill_hex(shape):
    try:
        if shape.fill.type == 1:  # MSO_FILL.SOLID
            return str(shape.fill.fore_color.rgb)
    except Exception:
        return None
    return None


def rects_overlap_area(a, b):
    l = max(a[0], b[0])
    r = min(a[0] + a[2], b[0] + b[2])
    t = max(a[1], b[1])
    bo = min(a[1] + a[3], b[1] + b[3])
    if r <= l or bo <= t:
        return 0.0
    return (r - l) * (bo - t)


def check_1_text_overflow(prs):
    """desc 텍스트박스 높이가 실제 필요한 wrap 줄 수보다 부족한지(잘림 위험)."""
    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            full_text = tf.text
            if not full_text or not full_text.strip():
                continue
            left = emu_to_in(shape.left)
            width = emu_to_in(shape.width)
            height = emu_to_in(shape.height)
            if left < SLIDE_W_IN * DESC_PANEL_LEFT_RATIO:
                continue
            wrap = tf.word_wrap
            if wrap is None:
                wrap = True

            para_info = []
            for para in tf.paragraphs:
                runs = para.runs
                if not runs:
                    para_text, pt, bold = '', 6.9, False
                else:
                    para_text = ''.join(r.text for r in runs)
                    pt, bold = None, False
                    for r in runs:
                        if r.font.size is not None:
                            pt = r.font.size.pt
                        if r.font.bold:
                            bold = True
                    if pt is None:
                        pt = 6.9
                ls = para.line_spacing
                if ls is not None and hasattr(ls, 'pt'):
                    line_h_pt = ls.pt
                elif isinstance(ls, (int, float)):
                    line_h_pt = ls * pt
                else:
                    line_h_pt = pt * 1.5
                para_info.append((para_text, pt, bold, line_h_pt))

            if not wrap:
                needed_h_in = sum(li[3] for li in para_info) / 72.0
            else:
                needed_h_in = 0.0
                for para_text, pt, bold, line_h_pt in para_info:
                    n_lines = 1 if para_text == '' else wrap_lines(para_text, pt, bold, width)
                    needed_h_in += n_lines * (line_h_pt / 72.0)

            overflow_in = needed_h_in - height
            if overflow_in > 0.08:
                findings.append({
                    'slide': si + 1, 'shape': shape.name,
                    'text': full_text.strip().replace('\n', ' | ')[:100],
                    'overflow_in': round(overflow_in, 2),
                })
    return findings


def check_2_desc_text_overlap(prs):
    """Description 패널 안에서 서로 다른(긴) 텍스트박스끼리 겹치는지."""
    findings = []
    for si, slide in enumerate(prs.slides):
        boxes = []
        for s in slide.shapes:
            if not (s.has_text_frame and s.text_frame.text.strip()):
                continue
            left = emu_to_in(s.left)
            if left < SLIDE_W_IN * DESC_PANEL_LEFT_RATIO:
                continue
            txt = s.text_frame.text.strip()
            if len(txt) < 15:  # 짧은 인디케이터 배지(예: "5-1")는 의도적으로 겹치므로 제외
                continue
            rect = (left, emu_to_in(s.top), emu_to_in(s.width), emu_to_in(s.height))
            boxes.append((rect, txt[:40]))
        for a in range(len(boxes)):
            for b in range(a + 1, len(boxes)):
                r1, t1 = boxes[a]
                r2, t2 = boxes[b]
                ov = rects_overlap_area(r1, r2)
                area1, area2 = r1[2] * r1[3], r2[2] * r2[3]
                if area1 <= 0 or area2 <= 0:
                    continue
                ratio = ov / min(area1, area2)
                if ratio > 0.05:
                    findings.append({'slide': si + 1, 'ratio': round(ratio * 100), 'a': t1, 'b': t2})
    return findings


def check_3_bg_covers_text(prs):
    """텍스트박스보다 나중에(위에) 그려지는 불투명 배경 도형이 텍스트를 가리는지."""
    findings = []
    for si, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        items = []
        for zi, shape in enumerate(shapes):
            left, top = emu_to_in(shape.left), emu_to_in(shape.top)
            w, h = emu_to_in(shape.width), emu_to_in(shape.height)
            if w <= 0 or h <= 0:
                continue
            has_text = bool(shape.has_text_frame and shape.text_frame.text.strip())
            fill_hex = get_fill_hex(shape) if not has_text else None
            items.append({
                'z': zi, 'rect': (left, top, w, h), 'has_text': has_text,
                'text': shape.text_frame.text.strip() if has_text else '',
                'is_opaque_shape': fill_hex is not None, 'left': left,
            })
        for i, it in enumerate(items):
            if not it['has_text'] or it['left'] < SLIDE_W_IN * DESC_PANEL_LEFT_RATIO:
                continue
            tr = it['rect']
            tr_area = tr[2] * tr[3]
            if tr_area <= 0:
                continue
            covered = 0.0
            for j in range(i + 1, len(items)):
                other = items[j]
                if not other['is_opaque_shape']:
                    continue
                covered += rects_overlap_area(tr, other['rect'])
            ratio = covered / tr_area
            if ratio > 0.3:
                findings.append({'slide': si + 1, 'ratio': round(ratio * 100), 'text': it['text'][:80]})
    return findings


def check_4_boundary(prs):
    """어떤 도형이든 슬라이드 하단/우측 경계를 벗어나는지."""
    findings = []
    for si, slide in enumerate(prs.slides):
        for s in slide.shapes:
            top, h = emu_to_in(s.top), emu_to_in(s.height)
            bottom = top + h
            if bottom > SLIDE_H_IN + 0.1:
                txt = s.text_frame.text[:40] if s.has_text_frame else s.name
                findings.append({'slide': si + 1, 'bottom': round(bottom, 2), 'text': txt})
    return findings


def check_5_bg_text_size_mismatch(prs):
    """desc-table 짝수행 배경(F8F9FA)이 자기 텍스트박스보다 작아서 텍스트가 삐져나오는지."""
    findings = []
    for si, slide in enumerate(prs.slides):
        text_boxes, bg_boxes = [], []
        for s in slide.shapes:
            left = emu_to_in(s.left)
            if left < SLIDE_W_IN * DESC_PANEL_LEFT_RATIO:
                continue
            rect = (left, emu_to_in(s.top), emu_to_in(s.width), emu_to_in(s.height))
            if s.has_text_frame and s.text_frame.text.strip() and len(s.text_frame.text.strip()) > 15:
                text_boxes.append((rect, s.text_frame.text.strip()[:40]))
            elif get_fill_hex(s) == 'F8F9FA':
                bg_boxes.append(rect)
        for trect, ttext in text_boxes:
            best, best_ov = None, 0
            for brect in bg_boxes:
                ov = rects_overlap_area(trect, brect)
                if ov > best_ov:
                    best_ov, best = ov, brect
            if best is None:
                continue
            text_bottom = trect[1] + trect[3]
            bg_bottom = best[1] + best[3]
            if text_bottom - bg_bottom > 0.03:
                findings.append({'slide': si + 1, 'excess': round(text_bottom - bg_bottom, 2), 'text': ttext})
    return findings


def main():
    path = sys.argv[1] if len(sys.argv) > 1 else 'downloaded.pptx'
    prs = Presentation(path)
    print(f"검증 대상: {path} ({len(prs.slides)} slides)\n")

    checks = [
        ("1. 텍스트 오버플로우(박스보다 필요한 줄 수가 많음 → 잘림 위험)", check_1_text_overflow),
        ("2. Description 텍스트박스끼리 겹침", check_2_desc_text_overlap),
        ("3. 배경 도형이 텍스트를 덮는 겹침", check_3_bg_covers_text),
        ("4. 슬라이드 경계(하단) 이탈", check_4_boundary),
        ("5. 배경 셀이 텍스트박스보다 작음(텍스트가 배경 밖으로 삐져나옴)", check_5_bg_text_size_mismatch),
    ]

    total_findings = 0
    for label, fn in checks:
        findings = fn(prs)
        total_findings += len(findings)
        status = "OK" if not findings else f"FAIL ({len(findings)}건)"
        print(f"[{status}] {label}")
        for f in findings[:15]:
            print(f"    {f}")
        if len(findings) > 15:
            print(f"    ... 외 {len(findings) - 15}건 더")
        print()

    print("=" * 60)
    if total_findings == 0:
        print("전체 통과 — 잘림/겹침/경계이탈 없음.")
    else:
        print(f"총 {total_findings}건 발견 — 실제 PPTX를 열어 해당 슬라이드 육안 확인 필요.")
    return 0 if total_findings == 0 else 1


if __name__ == '__main__':
    sys.exit(main())
